import sys, os, pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_flow_slides():
    pptx_path = 'PRESENTACION_PROYECTO.pptx'
    logo_png = 'Publica navojoa logotipo oficial.png'

    prs = pptx.Presentation(pptx_path)

    blank_slide_layout = prs.slide_layouts[6] # blank layout

    # Color Palette
    DARK_NAVY = RGBColor(15, 23, 42)
    EMERALD_GREEN = RGBColor(37, 211, 102)
    GOLD_ACCENT = RGBColor(217, 119, 6)
    LIGHT_BG = RGBColor(248, 250, 252)
    CARD_BG = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240)
    TEXT_MUTED = RGBColor(100, 116, 139)

    # ----------------------------------------------------
    # SLIDE 1: Flujo Completo de Captación y Atención 2026
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)

    # Header Title
    txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Flujo Operativo Integrado — Publica Navojoa"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = DARK_NAVY

    p2 = tf.add_paragraph()
    p2.text = "Sincronización 100% Automática: Grupo Facebook ➔ Registro Web ➔ Bot WhatsApp ➔ CRM en Vivo"
    p2.font.size = Pt(15)
    p2.font.color.rgb = EMERALD_GREEN

    # Add 4 Flow Step Cards
    steps = [
        ("1. Registro VIP Web", "publicanavojoa.com/registro", [
            "• Elección de Colonia/Zona en Navojoa.",
            "• Selector de Interés comercial:",
            "  - 🛍️ No, solo ofertas (Comprador VIP)",
            "  - 💼 Sí, me interesa anunciar (Prospecto)",
            "• Registro en tiempo real a Firebase.",
            "• Auto-redirección limpia a WhatsApp en 1.5s."
        ], RGBColor(239, 246, 255), RGBColor(59, 130, 246)),

        ("2. Bot WhatsApp Cloud", "+52 642 152 0280", [
            "• Reconoce al usuario por Nombre y Colonia.",
            "• Filtra usuarios no registrados.",
            "• Respuesta consultiva en ANUNCIAR:",
            "  - Asesor humano contacta en < 24 hrs.",
            "• Catálogo limpio y ágil en OFERTAS.",
            "• Despedida educada en ok, gracias, perfecto."
        ], RGBColor(240, 253, 244), EMERALD_GREEN),

        ("3. Panel CRM & Chat", "publicanavojoa.com/admin", [
            "• Pipeline Comercial de 5 Estados:",
            "  - 💼 Prospecto Anunciante",
            "  - 📞 Contacto Realizado",
            "  - 🟢 Cliente Activo / Cerrado",
            "  - 🔴 No Interesado",
            "  - 🛍️ Comprador VIP",
            "• Borrado de historial y sincronización directa."
        ], RGBColor(254, 243, 199), GOLD_ACCENT),

        ("4. Chat Multimedios", "Proxy Servidor Vercel", [
            "• Comunicación Bidireccional en Tiempo Real.",
            "• Envío/Recepción de fotos (PNG/JPG) a color.",
            "• Envío/Recepción de PDFs interactivos.",
            "• Proxy /api/download-media:",
            "  - Descarga directa sin error 401 de Meta."
        ], RGBColor(248, 250, 252), DARK_NAVY)
    ]

    left_margin = Inches(0.6)
    card_width = Inches(2.8)
    card_height = Inches(5.2)
    gap = Inches(0.25)

    for i, (title, subtitle, bullets, bg_color, accent_color) in enumerate(steps):
        x = left_margin + i * (card_width + gap)
        y = Inches(1.8)

        # Background Card Shape
        shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = accent_color
        shape.line.width = Pt(1.5)

        # Card Text Frame
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        # Title
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = accent_color

        # Subtitle
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.bold = True
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = DARK_NAVY
        p_sub.space_after = Pt(10)

        # Bullets
        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = b
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = DARK_NAVY
            p_b.space_after = Pt(4)

    # Insert logo at bottom right
    if os.path.exists(logo_png):
        slide1.shapes.add_picture(logo_png, Inches(10.2), Inches(0.4), width=Inches(2.5))


    # ----------------------------------------------------
    # SLIDE 2: Detalle Técnico y Módulos de Producción
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)

    # Header Title
    txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Arquitectura del Sistema y Módulos en Vivo"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = DARK_NAVY

    p2 = tf2.add_paragraph()
    p2.text = "Infraestructura desplegada y operando en producción en publicanavojoa.com"
    p2.font.size = Pt(15)
    p2.font.color.rgb = EMERALD_GREEN

    # 3 Large Feature Columns
    cols = [
        ("🌐 Portal Web & Tarifario", "publicanavojoa.com", [
            "• Página de Inicio (/): Media Kit y Tarifario Publicitario 2026.",
            "• Formulario (/registro): Formulario VIP con Firebase Firestore.",
            "• Privacidad (/privacidad): Políticas de privacidad y Opt-in Meta.",
            "• Vercel Rewrites: URLs limpias sin extensión .html."
        ], RGBColor(239, 246, 255), RGBColor(59, 130, 246)),

        ("💬 WhatsApp Cloud API & Bot", "+52 642 152 0280 (Línea Oficial)", [
            "• Webhook Serverless (/api/webhook): Validación de firmas Meta.",
            "• Lógica Consultiva: Asesoramiento personalizado en ANUNCIAR.",
            "• Sincronización Real-Time: Guarda automáticamente mensajes en Firebase.",
            "• Procesamiento de Medios: Decodificación de fotos y archivos PDF."
        ], RGBColor(240, 253, 244), EMERALD_GREEN),

        ("💼 Panel CRM & Proxy Media", "publicanavojoa.com/admin", [
            "• Gestión de Prospectos: Cambio de estado comercial en 1 clic.",
            "• Chat Multimedios: Envío/recepción de fotos y documentos.",
            "• Botón Borrar Chat: Limpieza individual de historial por contacto.",
            "• Proxy Media (/api/download-media): Descarga de archivos sin 401."
        ], RGBColor(254, 243, 199), GOLD_ACCENT)
    ]

    col_width = Inches(3.8)
    col_gap = Inches(0.3)
    col_left = Inches(0.7)

    for i, (title, subtitle, bullets, bg_color, accent_color) in enumerate(cols):
        x = col_left + i * (col_width + col_gap)
        y = Inches(1.8)

        shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_width, Inches(5.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = accent_color
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = accent_color

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.bold = True
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = DARK_NAVY
        p_sub.space_after = Pt(12)

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = b
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = DARK_NAVY
            p_b.space_after = Pt(6)

    if os.path.exists(logo_png):
        slide2.shapes.add_picture(logo_png, Inches(10.2), Inches(0.4), width=Inches(2.5))

    prs.save(pptx_path)
    print(f'SUCCESS: 2 flow slides added to {pptx_path}. Total slides: {len(prs.slides)}')

if __name__ == '__main__':
    add_flow_slides()
