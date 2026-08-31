import sys, os, pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def update_presentation():
    pptx_path = 'PRESENTACION_PROYECTO.pptx'
    logo_png = 'Publica navojoa logotipo oficial.png'

    prs = pptx.Presentation(pptx_path)

    # --- 1. Slide 1 (Portada Principal) ---
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if 'Red Publicitaria Digital Navojoa' in txt:
                shape.text_frame.text = 'PUBLICA NAVOJOA'
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(44)
                p.font.color.rgb = RGBColor(15, 23, 42)
            elif 'Publicidad Orgánica que Llega Directo' in txt:
                shape.text_frame.text = 'Privado, personalizado y directo.'
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(37, 211, 102)

    # Insert logo on Slide 1
    if os.path.exists(logo_png):
        s1.shapes.add_picture(logo_png, Inches(1.0), Inches(0.4), width=Inches(3.6))

    # --- 2. Slide 18 (Separador Parte 5) ---
    s18 = prs.slides[17]
    for shape in s18.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if 'Parte 5: Opciones de Nombre y Marca' in txt:
                shape.text_frame.text = 'Parte 5: Identidad de Marca Oficial'
            elif '3 Propuestas para Decidir en Equipo' in txt:
                shape.text_frame.text = 'Publica Navojoa — Logotipo y Slogan Aprobados'

    # --- 3. Slide 19 (Detalle de la Marca) ---
    s19 = prs.slides[18]
    for shape in s19.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if '3 Propuestas de Nombre Comercial' in txt:
                shape.text_frame.text = 'Identidad Oficial: Publica Navojoa'
            elif 'Vitrina Navojoa' in txt:
                shape.text_frame.text = 'Slogan Oficial'
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
            elif 'Evoca un escaparate digital' in txt:
                shape.text_frame.text = '“Privado, personalizado y directo.”\n\nResalta que el mensaje llega de forma 1 a 1, privada y directa al smartphone de cada persona en Navojoa.'
            elif 'Web: vitrinanavojoa.com' in txt:
                shape.text_frame.text = 'Difusión Directa al Celular'
            elif 'Navojoa Digital' in txt:
                shape.text_frame.text = 'Ecosistema Integrado'
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
            elif 'Posiciona la marca como LA plataforma' in txt:
                shape.text_frame.text = '• Grupo de Facebook (78.7K miembros)\n• Portal Web de Registro Opt-in\n• Bot de WhatsApp Automatizado\n• Envíos Masivos 100% Autorizados'
            elif 'Web: navojoadigital.com' in txt:
                shape.text_frame.text = 'Web: publicanavojoa.com'
            elif 'Nombre directo que invita a la acción' in txt:
                shape.text_frame.text = 'Nombre contundente, claro y enfocado en la difusión de todo tipo de anuncios (comerciales, políticos, eventos y comunitarios).'

    if os.path.exists(logo_png):
        s19.shapes.add_picture(logo_png, Inches(0.8), Inches(4.5), width=Inches(3.3))

    # --- 4. Slide 32 (Cronograma) ---
    s32 = prs.slides[31]
    for shape in s32.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if 'Definir nombre de marca y diseñar logo' in txt:
                shape.text_frame.text = txt.replace('Definir nombre de marca y diseñar logo.', '✅ Nombre de marca, slogan y logotipo oficial completados.')

    # --- 5. Slide 34 (Decisiones Pendientes) ---
    s34 = prs.slides[33]
    for shape in s34.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if 'ELEGIR EL NOMBRE DE LA MARCA' in txt:
                shape.text_frame.text = '✅ NOMBRE, SLOGAN Y LOGOTIPO DEFINIDOS: \"Publica Navojoa\" con slogan \"Privado, personalizado y directo.\" e identidad visual lista.'

    # --- 6. Slide 35 (Cierre / Portada Final) ---
    s35 = prs.slides[34]
    for shape in s35.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if 'Publicidad Orgánica. Sin Campañas Pagadas. Directo al Cliente.' in txt:
                shape.text_frame.text = 'PUBLICA NAVOJOA — Privado, personalizado y directo.'
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.color.rgb = RGBColor(37, 211, 102)

    if os.path.exists(logo_png):
        s35.shapes.add_picture(logo_png, Inches(4.5), Inches(0.8), width=Inches(4.3))

    prs.save(pptx_path)
    print('SUCCESS: PRESENTACION_PROYECTO.pptx updated.')

if __name__ == '__main__':
    update_presentation()
